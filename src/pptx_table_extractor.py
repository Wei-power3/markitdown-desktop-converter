"""
PPTX table extraction and conversion to Markdown tables
"""
import logging
from typing import List, Dict, Optional
from pathlib import Path

# Setup logging
logger = logging.getLogger(__name__)

# Try importing pptx
try:
    from pptx import Presentation
    from pptx.table import Table
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx not available. PPTX table extraction disabled.")


class PPTXTableExtractor:
    """
    Detects tables in PPTX and converts to Markdown table format.
    
    Fixes the issue where PPTX tables are extracted as flat bullet lists,
    losing column/row relationships.
    """
    
    def __init__(self):
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx is required for PPTX table extraction")
    
    def extract_tables_from_pptx(self, pptx_path: Path) -> List[Dict]:
        """
        Extract all tables from PPTX file.
        
        Args:
            pptx_path: Path to PPTX file
        
        Returns:
            List of table dictionaries with structure:
            {
                'slide_number': int,
                'slide_title': str,
                'rows': [[cell1, cell2, ...], ...],
                'row_count': int,
                'col_count': int,
                'markdown': str  # Pre-formatted markdown
            }
        """
        if not PPTX_AVAILABLE:
            return []
        
        try:
            prs = Presentation(str(pptx_path))
        except Exception as e:
            logger.error(f"Failed to open PPTX: {e}")
            return []
        
        tables = []
        
        for slide_idx, slide in enumerate(prs.slides):
            slide_num = slide_idx + 1
            slide_title = self._get_slide_title(slide)
            
            # Find table shapes
            for shape in slide.shapes:
                if shape.has_table:
                    try:
                        table_data = self._extract_table_data(shape.table)
                        
                        if not table_data:  # Skip empty tables
                            continue
                        
                        markdown_table = self._format_as_markdown(table_data)
                        
                        tables.append({
                            'slide_number': slide_num,
                            'slide_title': slide_title,
                            'rows': table_data,
                            'row_count': len(table_data),
                            'col_count': len(table_data[0]) if table_data else 0,
                            'markdown': markdown_table
                        })
                        
                        logger.info(
                            f"Extracted table from slide {slide_num} ({slide_title}): "
                            f"{len(table_data)} rows × {len(table_data[0])} cols"
                        )
                    except Exception as e:
                        logger.warning(f"Failed to extract table from slide {slide_num}: {e}")
                        continue
        
        return tables
    
    def _get_slide_title(self, slide) -> str:
        """
        Extract slide title (usually first text shape or title placeholder).
        
        Args:
            slide: python-pptx Slide object
        
        Returns:
            Slide title or "Untitled Slide"
        """
        # Try title placeholder first
        if slide.shapes.title:
            title = slide.shapes.title.text.strip()
            if title:
                return title
        
        # Fallback: find first text shape
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                return shape.text.strip()
        
        return "Untitled Slide"
    
    def _extract_table_data(self, table: Table) -> List[List[str]]:
        """
        Extract table data as 2D array.
        
        Args:
            table: python-pptx Table object
        
        Returns:
            [[cell1, cell2, ...], [cell1, cell2, ...], ...]
        """
        rows = []
        
        for row in table.rows:
            row_data = []
            for cell in row.cells:
                cell_text = cell.text.strip()
                
                # Clean cell text:
                # - Remove hard line breaks within cells
                # - Replace newlines with spaces
                # - Collapse multiple spaces
                cell_text = cell_text.replace('\n', ' ').replace('\r', '')
                cell_text = ' '.join(cell_text.split())  # Collapse spaces
                
                row_data.append(cell_text)
            
            rows.append(row_data)
        
        return rows
    
    def _format_as_markdown(self, table_data: List[List[str]]) -> str:
        """
        Convert 2D array to Markdown table.
        
        Format:
        | Header1 | Header2 | Header3 |
        |---------|---------|------|
        | Cell1   | Cell2   | Cell3   |
        
        Args:
            table_data: 2D list of cell values
        
        Returns:
            Markdown-formatted table string
        """
        if not table_data:
            return ""
        
        # Determine column widths for alignment
        num_cols = len(table_data[0])
        col_widths = [0] * num_cols
        
        for row in table_data:
            for i, cell in enumerate(row):
                col_widths[i] = max(col_widths[i], len(cell))
        
        # Minimum column width
        col_widths = [max(w, 3) for w in col_widths]
        
        lines = []
        
        # Header row (first row)
        header_cells = [
            cell.ljust(col_widths[i]) 
            for i, cell in enumerate(table_data[0])
        ]
        lines.append("| " + " | ".join(header_cells) + " |")
        
        # Separator row
        separators = ["-" * width for width in col_widths]
        lines.append("| " + " | ".join(separators) + " |")
        
        # Data rows
        for row in table_data[1:]:
            cells = [
                cell.ljust(col_widths[i]) 
                for i, cell in enumerate(row)
            ]
            lines.append("| " + " | ".join(cells) + " |")
        
        return "\n".join(lines)
    
    def format_tables_for_injection(self, tables: List[Dict]) -> Dict[int, str]:
        """
        Format tables for injection back into markdown by slide number.
        
        Args:
            tables: List of table dictionaries from extract_tables_from_pptx
        
        Returns:
            Dictionary mapping slide_number -> markdown table string
        """
        slide_tables = {}
        
        for table in tables:
            slide_num = table['slide_number']
            
            # Add table with context
            table_section = f"\n\n### Table ({table['row_count']}×{table['col_count']})\n\n"
            table_section += table['markdown']
            table_section += "\n\n"
            
            if slide_num in slide_tables:
                slide_tables[slide_num] += table_section
            else:
                slide_tables[slide_num] = table_section
        
        return slide_tables
    
    def get_extraction_report(self, tables: List[Dict]) -> Dict:
        """
        Generate a report of table extraction.
        
        Args:
            tables: List of extracted tables
        
        Returns:
            Dictionary with extraction statistics
        """
        if not tables:
            return {
                'total_tables': 0,
                'slides_with_tables': 0,
                'total_rows': 0,
                'total_cells': 0
            }
        
        slides_with_tables = set(t['slide_number'] for t in tables)
        total_rows = sum(t['row_count'] for t in tables)
        total_cells = sum(t['row_count'] * t['col_count'] for t in tables)
        
        return {
            'total_tables': len(tables),
            'slides_with_tables': len(slides_with_tables),
            'total_rows': total_rows,
            'total_cells': total_cells,
            'avg_rows_per_table': total_rows / len(tables) if tables else 0,
            'slide_numbers': sorted(slides_with_tables)
        }


# Example usage and testing
if __name__ == "__main__":
    import sys
    
    if len(sys.argv) < 2:
        print("Usage: python pptx_table_extractor.py <path_to_pptx>")
        sys.exit(1)
    
    pptx_path = Path(sys.argv[1])
    
    if not pptx_path.exists():
        print(f"Error: File not found: {pptx_path}")
        sys.exit(1)
    
    print(f"Extracting tables from: {pptx_path.name}\n")
    
    extractor = PPTXTableExtractor()
    tables = extractor.extract_tables_from_pptx(pptx_path)
    
    if not tables:
        print("No tables found in PPTX.")
    else:
        print(f"Found {len(tables)} table(s):\n")
        
        for i, table in enumerate(tables, 1):
            print(f"Table {i}:")
            print(f"  Slide: {table['slide_number']} ({table['slide_title']})")
            print(f"  Size: {table['row_count']} rows × {table['col_count']} cols")
            print(f"\nMarkdown output:")
            print(table['markdown'])
            print("\n" + "="*80 + "\n")
        
        # Statistics
        report = extractor.get_extraction_report(tables)
        print("\n=== Extraction Report ===")
        print(f"Total tables: {report['total_tables']}")
        print(f"Slides with tables: {report['slides_with_tables']}")
        print(f"Total rows: {report['total_rows']}")
        print(f"Total cells: {report['total_cells']}")
        print(f"Average rows per table: {report['avg_rows_per_table']:.1f}")
        print(f"Slide numbers: {report['slide_numbers']}")
