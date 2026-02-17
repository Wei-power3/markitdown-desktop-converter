"""
Document conversion logic using MarkItDown with enhanced quality
"""
from pathlib import Path
from typing import Tuple, Optional, List, Dict
import re
from markitdown import MarkItDown
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.units import inch
from pptx import Presentation
import io
import logging
from config import CONVERT_PPTX_TO_PDF, CONVERT_PPTX_DIRECT
from text_cleaner import MarkdownCleaner
from table_extractor import TableExtractor
from pptx_table_extractor import PPTXTableExtractor

# Setup logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class DocumentConverter:
    """
    Enhanced document conversion with:
    - Text cleaning (fixes ligatures, spacing, encoding)
    - PPTX run-on word repair (fixes "withabusiness" → "with a business")
    - PPTX table structure preservation (converts tables to proper markdown)
    - Structured table extraction
    - Dual PowerPoint conversion
    - Word document support
    - Excel spreadsheet support
    """
    
    def __init__(self):
        self.md = MarkItDown()
        self.text_cleaner = MarkdownCleaner()
        self.table_extractor = TableExtractor(min_accuracy=0.5)  # For PDF tables
        
        # PPTX table extractor
        try:
            self.pptx_table_extractor = PPTXTableExtractor()
        except ImportError:
            self.pptx_table_extractor = None
            logger.warning("PPTX table extraction disabled (python-pptx not available)")
    
    def convert_file(self, file_path: Path) -> Tuple[str, Optional[str]]:
        """
        Convert file to Markdown with enhanced quality.
        
        Args:
            file_path: Path to file to convert
        
        Returns:
            Tuple of (markdown_content, error_message)
            error_message is None if successful
        """
        try:
            extension = file_path.suffix.lower()
            
            if extension == '.pdf':
                return self._convert_pdf(file_path)
            elif extension in ['.pptx', '.ppt']:
                return self._convert_powerpoint(file_path)
            elif extension in ['.docx', '.doc']:
                return self._convert_word(file_path)
            elif extension in ['.xlsx', '.xls']:
                return self._convert_excel(file_path)
            else:
                return "", f"Unsupported file type: {extension}"
        
        except Exception as e:
            logger.error(f"Conversion error: {str(e)}", exc_info=True)
            return "", f"Conversion error: {str(e)}"
    
    def _convert_pdf(self, file_path: Path) -> Tuple[str, None]:
        """
        Convert PDF to Markdown with enhanced quality.
        
        Process:
        1. Extract base text using MarkItDown
        2. Clean text artifacts (ligatures, spacing, etc.)
        3. Extract structured tables
        4. Combine all content
        """
        logger.info(f"Converting PDF: {file_path.name}")
        
        # Step 1: Base conversion
        result = self.md.convert(str(file_path))
        base_content = result.text_content
        
        # Step 2: Clean text (PDF format)
        cleaned_content = self.text_cleaner.clean(base_content, source_format='pdf')
        
        # Log cleaning statistics
        report = self.text_cleaner.get_cleaning_report(base_content, cleaned_content, source_format='pdf')
        logger.info(f"Text cleaning: {report['encoding_fixes']} encoding fixes, "
                   f"{report['hyphen_fixes']} hyphen fixes, "
                   f"{report['medical_term_fixes']} medical term fixes")
        
        # Step 3: Extract structured tables
        try:
            tables = self.table_extractor.extract_tables(file_path)
            if tables:
                logger.info(f"Extracted {len(tables)} structured table(s)")
                table_section = self.table_extractor.format_tables_for_markdown(tables)
                cleaned_content += table_section
        except Exception as e:
            logger.warning(f"Table extraction failed: {e}")
            # Continue without tables - not critical
        
        return cleaned_content, None
    
    def _convert_powerpoint(self, file_path: Path) -> Tuple[str, None]:
        """
        Convert PowerPoint to Markdown with table structure preservation.
        
        Process:
        1. Extract tables FIRST using python-pptx
        2. Convert PPTX to markdown using MarkItDown
        3. Apply PPTX-specific text cleaning (run-on words)
        4. Inject properly formatted tables back into content
        
        Implements both direct conversion and PDF pathway.
        """
        logger.info(f"Converting PowerPoint: {file_path.name}")
        results = []
        
        # Step 1: Extract tables BEFORE text conversion
        pptx_tables = []
        if self.pptx_table_extractor:
            try:
                pptx_tables = self.pptx_table_extractor.extract_tables_from_pptx(file_path)
                if pptx_tables:
                    report = self.pptx_table_extractor.get_extraction_report(pptx_tables)
                    logger.info(
                        f"PPTX table extraction: {report['total_tables']} tables found "
                        f"on slides {report['slide_numbers']}"
                    )
            except Exception as e:
                logger.warning(f"PPTX table extraction failed: {e}")
        
        # Method 1: Direct PPTX → Markdown
        if CONVERT_PPTX_DIRECT:
            try:
                result = self.md.convert(str(file_path))
                direct_content = result.text_content
                
                # Store original for comparison
                original_content = direct_content
                
                # Step 2: Clean the content with PPTX-specific fixes
                direct_content = self.text_cleaner.clean(direct_content, source_format='pptx')
                
                # Step 3: Inject tables if we have them
                if pptx_tables:
                    direct_content = self._inject_pptx_tables(direct_content, pptx_tables)
                    logger.info(f"Injected {len(pptx_tables)} table(s) into markdown")
                
                # Log PPTX-specific statistics
                report = self.text_cleaner.get_cleaning_report(
                    original_content, 
                    direct_content, 
                    source_format='pptx'
                )
                logger.info(
                    f"PPTX text cleaning: "
                    f"{report.get('pptx_contraction_fixes', 0)} contraction fixes, "
                    f"{report.get('pptx_run_on_fixes', 0)} run-on word fixes, "
                    f"token increase: +{report.get('pptx_token_delta', 0)}"
                )
                
                results.append("# Direct PPTX Conversion\n\n")
                results.append(direct_content)
                logger.info("Direct PPTX conversion completed with quality fixes")
            except Exception as e:
                logger.warning(f"Direct conversion failed: {str(e)}")
                results.append(f"Direct conversion failed: {str(e)}\n\n")
        
        # Method 2: PPTX → PDF → Markdown
        if CONVERT_PPTX_TO_PDF:
            try:
                pdf_content = self._pptx_to_pdf(file_path)
                pdf_path = file_path.with_suffix('.pdf')
                pdf_path.write_bytes(pdf_content)
                
                result = self.md.convert(str(pdf_path))
                pdf_pathway_content = result.text_content
                
                # Clean the content (PPTX-originated, so use pptx format)
                pdf_pathway_content = self.text_cleaner.clean(pdf_pathway_content, source_format='pptx')
                
                # Inject tables
                if pptx_tables:
                    pdf_pathway_content = self._inject_pptx_tables(pdf_pathway_content, pptx_tables)
                
                results.append("\n\n---\n\n# PDF Pathway Conversion\n\n")
                results.append(pdf_pathway_content)
                
                # Clean up temporary PDF
                pdf_path.unlink()
                logger.info("PDF pathway conversion completed")
            except Exception as e:
                logger.warning(f"PDF pathway failed: {str(e)}")
                results.append(f"\n\nPDF pathway failed: {str(e)}")
        
        final_content = "".join(results)
        return final_content, None
    
    def _inject_pptx_tables(self, content: str, tables: List[Dict]) -> str:
        """
        Inject properly formatted tables into markdown content.
        
        Strategy:
        - Find slide comment markers (e.g., <!-- Slide number: 3 -->)
        - Insert table after the slide marker
        - If MarkItDown already created a table, enhance it
        - If no table exists, add our extracted table
        
        Args:
            content: Markdown content from MarkItDown
            tables: List of table dicts from PPTXTableExtractor
        
        Returns:
            Enhanced markdown with proper tables
        """
        if not tables:
            return content
        
        # Create slide number to tables mapping
        slide_tables = self.pptx_table_extractor.format_tables_for_injection(tables)
        
        # Process content line by line
        lines = content.split('\n')
        enhanced_lines = []
        current_slide = None
        table_injected = set()  # Track which slides have tables injected
        
        for i, line in enumerate(lines):
            enhanced_lines.append(line)
            
            # Detect slide markers: <!-- Slide number: X -->
            slide_match = re.match(r'<!--\s*Slide number:\s*(\d+)\s*-->', line)
            if slide_match:
                current_slide = int(slide_match.group(1))
                
                # If this slide has tables and we haven't injected yet
                if current_slide in slide_tables and current_slide not in table_injected:
                    # Look ahead to see if MarkItDown already created a table
                    has_table_ahead = False
                    for j in range(i+1, min(i+20, len(lines))):
                        if '|' in lines[j] and '---' in lines[j+1] if j+1 < len(lines) else False:
                            has_table_ahead = True
                            break
                    
                    # If no table ahead, inject ours
                    if not has_table_ahead:
                        enhanced_lines.append(slide_tables[current_slide])
                        table_injected.add(current_slide)
                        logger.debug(f"Injected table for slide {current_slide}")
        
        return '\n'.join(enhanced_lines)
    
    def _convert_word(self, file_path: Path) -> Tuple[str, None]:
        """
        Convert Word document (.docx, .doc) to Markdown.
        
        MarkItDown natively supports Word documents:
        - Extracts text content
        - Preserves basic structure
        - Handles formatting
        
        Text cleaning is applied for quality.
        """
        logger.info(f"Converting Word document: {file_path.name}")
        
        try:
            # MarkItDown handles Word natively
            result = self.md.convert(str(file_path))
            content = result.text_content
            
            # Apply text cleaning for quality
            cleaned_content = self.text_cleaner.clean(content, source_format='docx')
            
            # Log cleaning statistics
            report = self.text_cleaner.get_cleaning_report(content, cleaned_content, source_format='docx')
            logger.info(f"Word doc cleaned: {report['encoding_fixes']} encoding fixes, "
                       f"{report['hyphen_fixes']} hyphen fixes")
            
            return cleaned_content, None
            
        except Exception as e:
            logger.error(f"Word conversion failed: {str(e)}")
            return "", f"Word conversion error: {str(e)}"
    
    def _convert_excel(self, file_path: Path) -> Tuple[str, None]:
        """
        Convert Excel spreadsheet (.xlsx, .xls) to Markdown.
        
        MarkItDown natively supports Excel:
        - Converts sheets to markdown tables
        - Preserves cell data
        - Handles multiple sheets
        
        Text cleaning is applied for quality.
        """
        logger.info(f"Converting Excel spreadsheet: {file_path.name}")
        
        try:
            # MarkItDown handles Excel natively
            # It converts each sheet to a markdown table
            result = self.md.convert(str(file_path))
            content = result.text_content
            
            # Apply text cleaning for quality
            # (handles any text artifacts in cell values)
            cleaned_content = self.text_cleaner.clean(content, source_format='xlsx')
            
            # Log info
            num_lines = len(cleaned_content.split('\n'))
            logger.info(f"Excel converted: {num_lines} lines of markdown")
            
            return cleaned_content, None
            
        except Exception as e:
            logger.error(f"Excel conversion failed: {str(e)}")
            return "", f"Excel conversion error: {str(e)}"
    
    def _pptx_to_pdf(self, pptx_path: Path) -> bytes:
        """
        Convert PPTX to PDF using reportlab.
        Returns PDF as bytes.
        """
        prs = Presentation(str(pptx_path))
        
        # Create PDF in memory
        pdf_buffer = io.BytesIO()
        doc = SimpleDocTemplate(pdf_buffer, pagesize=letter)
        styles = getSampleStyleSheet()
        story = []
        
        for i, slide in enumerate(prs.slides):
            # Add slide number
            story.append(Paragraph(f"<b>Slide {i+1}</b>", styles['Heading1']))
            story.append(Spacer(1, 0.2*inch))
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    # Clean and format text
                    text = shape.text.strip()
                    if text:
                        story.append(Paragraph(text, styles['Normal']))
                        story.append(Spacer(1, 0.1*inch))
            
            # Add spacing between slides
            story.append(Spacer(1, 0.3*inch))
        
        doc.build(story)
        return pdf_buffer.getvalue()
